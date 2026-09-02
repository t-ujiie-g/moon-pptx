#!/usr/bin/env python3
"""python-pptx side of the cross-library benchmark.

Same workload as `tools/bench/moonbit` and `bench_pptxgenjs.js`: N slides,
one text box each, serialised to an in-memory buffer. Nothing touches the
disk — moon-pptx has no file I/O by design (ADR-002), so writing a file
would time the host, not the library.

Prints the byte count so the work cannot be optimised away.

    bench_pptx.py <slide-count>
"""
import io
import sys

from pptx import Presentation
from pptx.util import Emu

X, Y, CX, CY = Emu(457200), Emu(457200), Emu(8229600), Emu(914400)


def build_and_save(n: int) -> int:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for i in range(n):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(X, Y, CX, CY)
        box.text_frame.text = f"Slide {i + 1}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getbuffer().nbytes


if __name__ == "__main__":
    print(build_and_save(int(sys.argv[1])))
